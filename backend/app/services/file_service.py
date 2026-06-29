"""
Memoria AI — File Service

Handles extraction of text content from multiple file formats:
- PDFs (via pdfplumber)
- PPTs/PPTXs (via python-pptx)
- Images/Handwritten notes (via Claude Vision API OCR)
- Audio recordings/files (via Groq Whisper API)
"""

import io
import base64
import httpx
import pdfplumber
from pptx import Presentation
from anthropic import Anthropic
from app.config import get_settings


def extract_pdf_text(file_bytes: bytes) -> str:
    """Extract text from PDF pages using pdfplumber."""
    text_content = []
    with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
        for i, page in enumerate(pdf.pages):
            page_text = page.extract_text()
            if page_text:
                text_content.append(f"--- Page {i + 1} ---\n{page_text}")
    return "\n\n".join(text_content)


def extract_pptx_text(file_bytes: bytes) -> str:
    """Extract text content from PowerPoint presentation slides using python-pptx."""
    prs = Presentation(io.BytesIO(file_bytes))
    text_content = []
    for slide_idx, slide in enumerate(prs.slides):
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())
        if slide_text:
            text_content.append(f"--- Slide {slide_idx + 1} ---\n" + "\n".join(slide_text))
    return "\n\n".join(text_content)


def extract_image_text(file_bytes: bytes, api_key: str, mime_type: str = "image/png") -> str:
    """Extract printed or handwritten text from images using Anthropic Claude Vision."""
    if not api_key or api_key.strip().lower() in ("", "your_key", "your-key", "sk-ant-your-key-here"):
        # Fallback to backend config key
        settings = get_settings()
        api_key = settings.anthropic_api_key

    if not api_key or api_key.strip().lower() in ("", "your_key", "your-key", "sk-ant-your-key-here"):
        raise ValueError("Anthropic API key is missing or not configured. Cannot perform image OCR.")

    # Encode image bytes to base64
    base64_image = base64.b64encode(file_bytes).decode("utf-8")
    client = Anthropic(api_key=api_key)

    response = client.messages.create(
        model="claude-3-5-sonnet-20241022",
        max_tokens=2048,
        messages=[
            {
                "role": "user",
                "content": [
                    {
                        "type": "image",
                        "source": {
                            "type": "base64",
                            "media_type": mime_type,
                            "data": base64_image,
                        },
                    },
                    {
                        "type": "text",
                        "text": (
                            "Transcribe all text from this image exactly. If there are handwritten notes or drawings, "
                            "transcribe the textual components cleanly. Do not summarize or explain the image. "
                            "Output only the raw transcribed text without any introductory or concluding remarks."
                        ),
                    }
                ],
            }
        ],
    )
    return response.content[0].text


def transcribe_audio(file_bytes: bytes, filename: str, groq_api_key: str) -> str:
    """Transcribe audio files using Groq Whisper API (whisper-large-v3)."""
    if not groq_api_key or groq_api_key.strip().lower() in ("", "your_key", "your-key", "gsk_..."):
        settings = get_settings()
        groq_api_key = settings.groq_api_key

    if not groq_api_key or groq_api_key.strip().lower() in ("", "your_key", "your-key", "gsk_..."):
        raise ValueError("Groq API key is missing or not configured. Cannot transcribe audio.")

    # Map extension to MIME type
    ext = filename.split(".")[-1].lower()
    mime_type = "application/octet-stream"
    if ext == "mp3":
        mime_type = "audio/mpeg"
    elif ext == "wav":
        mime_type = "audio/wav"
    elif ext in ("m4a", "x-m4a"):
        mime_type = "audio/mp4"
    elif ext == "ogg":
        mime_type = "audio/ogg"
    elif ext == "webm":
        mime_type = "audio/webm"

    url = "https://api.groq.com/openai/v1/audio/transcriptions"
    headers = {
        "Authorization": f"Bearer {groq_api_key}",
    }
    files = {
        "file": (filename, file_bytes, mime_type),
    }
    data = {
        "model": "whisper-large-v3",
        "response_format": "json",
    }

    # Using httpx to upload the file to Groq
    with httpx.Client(timeout=120.0) as client:
        resp = client.post(url, headers=headers, files=files, data=data)
        if resp.status_code != 200:
            raise RuntimeError(f"Groq Whisper transcription failed: {resp.text}")
        return resp.json().get("text", "")
